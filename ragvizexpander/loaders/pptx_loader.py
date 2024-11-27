import io

import tempfile

from typing import List
import re
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from unstructured.partition.pptx import partition_pptx
from llama_index.core import SimpleDirectoryReader
from docling.document_converter import DocumentConverter

from .base import DocumentLoader, LoaderStrategy


class PptxNativeStrategy(LoaderStrategy):
    """Native PPTX loading strategy using python-pptx"""

    def _load_textframed_shapes(self, shapes):
        for shape in shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_shape = shape
                    for shape in self._load_textframed_shapes(group_shape.shapes):
                        yield shape
                    continue
            except:
                continue

            if shape.has_text_frame:
                yield shape

    def _load_text(self, shapes, title_text) -> List[str]:
        all_text = []
        for shape in self._load_textframed_shapes(shapes):
            text = re.sub("\s{2,}", " ", shape.text).strip()
            if text and text != title_text:
                all_text.append(text)
        return all_text

    def _load_single_table(self, table) -> List[List[str]]:
        n_row = len(table.rows)
        n_col = len(table.columns)
        arrays = [["" for _ in range(n_row)] for _ in range(n_col)]

        for i in range(n_row):
            for j in range(n_col):
                cell = table.cell(i, j).text_frame.paragraphs
                cell_text = "".join([run.text.strip() for para in cell for run in para.runs])
                cell_text = re.sub("\n", "<br />", cell_text).strip()
                arrays[j][i] = cell_text
        return arrays

    def load(self, file_path: str) -> List[str]:
        presentation = Presentation(file_path)
        texts = []
        
        for slide in presentation.slides:
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    texts.append(title_text)
            
            # Extract text from shapes
            texts.extend(self._load_text(slide.shapes, title_text))
            
            # Extract text from tables
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = self._load_single_table(shape.table)
                    texts.extend([cell for row in table_data for cell in row if cell.strip()])
        
        return texts


class PptxUnstructuredStrategy(LoaderStrategy):
    """Unstructured library loading strategy for PPTX"""

    def load(self, file_path: str) -> List[str]:
        elements = partition_pptx(file_path)
        return ["\n".join([ele.text.strip() for ele in elements])]


class PptxLlamaIndexStrategy(LoaderStrategy):
    """LlamaIndex loading strategy for PPTX"""

    def load(self, file_path: str) -> List[str]:
        reader = SimpleDirectoryReader(input_files=[file_path])
        document = reader.load_data()[0]
        return [document.text.strip()]


class PptxDoclingStrategy(LoaderStrategy):
    """docling loading strategy for PPTX"""

    def load(self, file_path) -> List[str]:
        converter = DocumentConverter()
        if isinstance(file_path, io.BytesIO):
            with tempfile.NamedTemporaryFile(suffix=".docx") as temp:
                temp.write(file_path.getbuffer())
                temp.seek(0)
                result = converter.convert(temp.name)
        else:
            result = converter.convert(file_path)
        return [result.document.export_to_markdown()]


class PptxLoader(DocumentLoader):
    """PPTX document loader with multiple loading strategies"""

    def __init__(self, strategy: LoaderStrategy = None):
        strategy = strategy or PptxNativeStrategy()
        super().__init__(strategy)

    def load_data(self, file_path: str) -> List[str]:
        if not self.validate_file(file_path):
            raise ValueError(f"Invalid or non-existent file: {file_path}")
        return self._strategy.load(file_path)

    def supported_extensions(self) -> List[str]:
        return ['.pptx']


# Register valid strategies for PptxLoader
PptxLoader.register_strategies([
    PptxNativeStrategy,
    PptxUnstructuredStrategy,
    PptxLlamaIndexStrategy,
    PptxDoclingStrategy,
])
